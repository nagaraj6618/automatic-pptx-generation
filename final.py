import openpyxl
from pptx import Presentation
from PIL import Image
import os
from pathlib import Path
from pathlib import Path
import openpyxl

def is_file_present(directory, filename):
   directory_path = Path(directory)
   if not directory_path.is_dir():
      print(f'The specified directory "{directory}" does not exist.')
      return False

   file_path = directory_path / filename
   return file_path.exists()

def log_missing_files(directory, filenames, log_filename):
   log_directory = Path(r'D:\graduation\program\excel')
   log_filename=log_filename+'-missing.xlsx'

   log_file_path = log_directory / log_filename
   log_directory.mkdir(parents=True, exist_ok=True)

   if not log_file_path.is_file():
      workbook = openpyxl.Workbook()
      sheet = workbook.active
      sheet.append(['Missing Files'])
      workbook.save(log_file_path)
    

   workbook = openpyxl.load_workbook(log_file_path)
   sheet = workbook.active
    
   if not is_file_present(directory, filenames):
      sheet.append([filenames])


   workbook.save(log_file_path)



def _add_image(slide, placeholder_id, image_url):
   placeholder = slide.placeholders[placeholder_id]
   if os.path.exists(image_url):
      im = Image.open(image_url)
      width, height = im.size

      placeholder.width = width
      placeholder.height = height
      placeholder = placeholder.insert_picture(image_url)

      image_ratio = width / height
      placeholder_ratio = placeholder.width / placeholder.height
      ratio_difference = placeholder_ratio - image_ratio

      if ratio_difference > 0:
         difference_on_each_side = ratio_difference / 2
         placeholder.crop_left = -difference_on_each_side
         placeholder.crop_right = -difference_on_each_side
      else:
         difference_on_each_side = -ratio_difference / 2
         placeholder.crop_bottom = -difference_on_each_side
         placeholder.crop_top = -difference_on_each_side



wb = openpyxl.load_workbook(r"D:\graduation\excelfile\ECE.xlsx")


sheet = wb.worksheets[0]

prs = Presentation("D:\graduation\IT.pptx")

rows=sheet.max_row
imagePath = r"D:\graduation\image\graduate-image\ece\\"
count=0
for i in range(3,rows+1):
   directory_path = imagePath
   file_name = str(sheet.cell(i,2).value)
 
   # print("hii")
   if (is_file_present(directory_path, file_name+'.jpg') or is_file_present(directory_path, file_name+' .jpg') or is_file_present(directory_path, file_name+'.png')):

      layout=prs.slide_layouts[8]
      slide=prs.slides.add_slide(layout)
      title= slide.placeholders[0].text=str(sheet.cell(i,2).value)+'\n'+str(sheet.cell(i,3).value)
      sub= slide.placeholders[2].text=str(sheet.cell(i,4).value)
      # sub1= slide.shapes.placeholders[2].text=str(sheet.cell(i,6).value)
   
      if is_file_present(directory_path, file_name+' .jpg'):

         _add_image(slide,1,imagePath + str(sheet.cell(i,2).value)+" .jpg")
      _add_image(slide,1,imagePath + str(sheet.cell(i,2).value)+".jpg")
      print('success..')
   else:
      count+=1
      log_missing_files(directory_path, file_name,str(sheet.cell(1,1).value))
      print(f'The file {file_name} is not present in the directory.')
prs.save("D:\graduation\program\PPT\\"+str(sheet.cell(1,1).value)+'.pptx')
os.startfile("D:\graduation\program\PPT\\"+str(sheet.cell(1,1).value)+'.pptx')
print(count)